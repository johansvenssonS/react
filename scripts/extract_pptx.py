#!/usr/bin/env python3
"""Extract slide text and speaker notes from a .pptx into a markdown file
under slides/notes/, so Claude can read course slides without needing
binary pptx support."""
import sys
from pathlib import Path

from pptx import Presentation


def extract(pptx_path: Path, out_path: Path) -> None:
    prs = Presentation(pptx_path)
    lines = [f"# {pptx_path.stem}", ""]
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                lines.append(f"> Speaker notes: {notes}")
        lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    if len(sys.argv) != 2:
        print("usage: extract_pptx.py <path-to.pptx>", file=sys.stderr)
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    out_dir = Path(__file__).resolve().parent.parent / "slides" / "notes"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{pptx_path.stem}.md"
    extract(pptx_path, out_path)
    print(out_path)


if __name__ == "__main__":
    main()
